"""
Generate the UniVerse third-year PROJECT defense slide deck (.pptx).

Follows the department "Third Year Defense - Project Slide Show Template":
 1. Project Title          2. Presentation Outline    3. Team Overview
 4. Project Overview/Goal  5. Functional Requirements 6. Use Case Diagram
 7. Non-Functional Req.    8. Technical Diagram: ER   9. Tech Diagram: DFD L0
10. Tech Diagram: DFD L1  11. Tech Diagram: Activity 12. Project Demonstration
13. Discussion/Conclusion 14. Demo

The template allows several consecutive slides for Functional Requirements,
Technical Diagrams and the Demo, so each diagram gets its own labeled slide.
All diagram + demo areas are intentionally left as blank placeholder boxes
(filled in manually by the team).

Run:  python tool/make_defense_pptx.py
Out:  docs/UniVerse_Project_Defense.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ----------------------------------------------------------------------------
# Palette (UniVerse brand)
# ----------------------------------------------------------------------------
ORANGE      = RGBColor(0xFF, 0x7A, 0x00)
ORANGE_DARK = RGBColor(0xE6, 0x6A, 0x00)
INK         = RGBColor(0x1A, 0x1A, 0x1C)   # near-black headings
SLATE       = RGBColor(0x2B, 0x32, 0x40)   # dark slate band
BODY        = RGBColor(0x33, 0x37, 0x3D)   # body text
MUTED       = RGBColor(0x6E, 0x72, 0x78)   # captions
LINE        = RGBColor(0xDD, 0xDF, 0xE3)   # hairlines
PH_BG       = RGBColor(0xF4, 0xF5, 0xF7)   # placeholder fill
PH_BORDER   = RGBColor(0xC4, 0xC8, 0xCE)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

FONT       = "Calibri"
FONT_LIGHT = "Calibri Light"
FONT_MONO  = "Consolas"

# 16:9
EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------------
def _no_line(shape):
    shape.line.fill.background()


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    _no_line(shape)


def add_rect(slide, l, t, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    _solid(sp, color)
    sp.shadow.inherit = False
    return sp


def add_text(slide, l, t, w, h, runs, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=0,
             wrap=True):
    """runs: list of paragraphs; each paragraph a list of
    (text, size, color, bold, font, italic)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for (text, size, color, bold, font, italic) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = font
    return tb


def R(text, size, color=BODY, bold=False, font=FONT, italic=False):
    return (text, size, color, bold, font, italic)


def bg(slide, color=WHITE):
    add_rect(slide, 0, 0, EMU_W, EMU_H, color)


def footer(slide, n):
    add_rect(slide, Inches(0.6), Inches(7.02), Inches(12.13), Pt(0.75), LINE)
    add_text(slide, Inches(0.6), Inches(7.08), Inches(9), Inches(0.3),
             [[R("UniVerse - A Campus Companion  ·  Team Sherlocked, CSE, Leading University",
                 9, MUTED)]], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(11.9), Inches(7.08), Inches(0.83), Inches(0.3),
             [[R(str(n), 9, MUTED, bold=True)]],
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def header(slide, title, kicker=None):
    """Standard content-slide header: kicker + title + accent rule."""
    title_top = Inches(0.88) if kicker else Inches(0.55)
    rule_top = Inches(1.5) if kicker else Inches(1.18)
    if kicker:
        add_text(slide, Inches(0.7), Inches(0.55), Inches(11), Inches(0.3),
                 [[R(kicker.upper(), 11, ORANGE, bold=True)]])
    add_text(slide, Inches(0.7), title_top, Inches(12), Inches(0.7),
             [[R(title, 28, INK, bold=True, font=FONT_LIGHT)]])
    add_rect(slide, Inches(0.72), rule_top, Inches(0.62), Pt(3.2), ORANGE)


def bullets(slide, l, t, w, h, items, *, size=15, gap=10, sub_size=12.5,
            anchor=MSO_ANCHOR.TOP):
    """items: list of (text, level). level 0 = orange dot, 1 = dash sub."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    for text, lvl in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = 1.06
        p.space_after = Pt(gap if lvl == 0 else gap * 0.45)
        p.space_before = Pt(0)
        if lvl == 0:
            dot = p.add_run()
            dot.text = "●  "
            dot.font.size = Pt(9)
            dot.font.color.rgb = ORANGE
            dot.font.name = FONT
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = BODY
            r.font.name = FONT
        else:
            p.level = 1
            dash = p.add_run()
            dash.text = "–  "
            dash.font.size = Pt(sub_size)
            dash.font.color.rgb = MUTED
            dash.font.name = FONT
            r = p.add_run()
            r.text = text
            r.font.size = Pt(sub_size)
            r.font.color.rgb = MUTED
            r.font.name = FONT
    return tb


def grouped(slide, l, t, w, h, groups, *, hsize=13.5, bsize=12, gap=5,
            group_gap=11):
    """groups: list of (header, [items]) -> bold header + orange-dot bullets."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    for head, items in groups:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.space_before = Pt(0 if first else group_gap)
        p.space_after = Pt(4)
        p.line_spacing = 1.0
        first = False
        r = p.add_run()
        r.text = head
        r.font.size = Pt(hsize)
        r.font.bold = True
        r.font.color.rgb = INK
        r.font.name = FONT
        for it in items:
            pp = tf.add_paragraph()
            pp.space_before = Pt(0)
            pp.space_after = Pt(gap)
            pp.line_spacing = 1.04
            dot = pp.add_run()
            dot.text = "●  "
            dot.font.size = Pt(8)
            dot.font.color.rgb = ORANGE
            dot.font.name = FONT
            rr = pp.add_run()
            rr.text = it
            rr.font.size = Pt(bsize)
            rr.font.color.rgb = BODY
            rr.font.name = FONT
    return tb


def placeholder(slide, l, t, w, h, label, sub=""):
    box = add_rect(slide, l, t, w, h, PH_BG, MSO_SHAPE.ROUNDED_RECTANGLE)
    box.line.color.rgb = PH_BORDER
    box.line.width = Pt(1.25)
    box.line.dash_style = 2  # dashed
    runs = [[R(label, 14, MUTED, bold=True)]]
    if sub:
        runs.append([R(sub, 11, MUTED, italic=True)])
    add_text(slide, l, t, w, h, runs,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=3)


def caption(slide, text):
    """One-line descriptor under the title for diagram slides."""
    add_text(slide, Inches(0.7), Inches(1.62), Inches(11.9), Inches(0.3),
             [[R(text, 12.5, BODY)]])


def new_slide():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    return s


# ============================================================================
# SLIDE 1 — PROJECT TITLE
# ============================================================================
s = new_slide()
add_rect(s, 0, 0, Inches(0.28), EMU_H, ORANGE)
add_rect(s, 0, Inches(6.55), EMU_W, Inches(0.95), SLATE)

add_text(s, Inches(1.0), Inches(0.7), Inches(11), Inches(0.35),
         [[R("THIRD YEAR PROJECT  ·  DEFENSE PRESENTATION", 12, ORANGE, bold=True)]])

add_text(s, Inches(0.95), Inches(1.55), Inches(11.4), Inches(1.7),
         [[R("UniVerse - A Campus Companion", 42, INK, bold=True, font=FONT_LIGHT)]],
         line_spacing=1.0)

add_text(s, Inches(1.0), Inches(2.85), Inches(11.2), Inches(0.9),
         [[R("An Automatic Department Timetable Generator (OR-Tools CP-SAT) "
             "with Real-Time Campus Services", 18, BODY)]],
         line_spacing=1.1)

add_rect(s, Inches(1.02), Inches(3.85), Inches(0.85), Pt(3.5), ORANGE)

team = [
    ("Fahmid Alam", "0182320012101309"),
    ("Swadheen Islam Robi", "0182320012101278"),
    ("Shahriar Rashid Ratul", "0182320012101276"),
]
col_w = Inches(3.75)
x = Inches(1.0)
for name, sid in team:
    add_text(s, x, Inches(4.25), col_w, Inches(0.6),
             [[R(name, 15, INK, bold=True)],
              [R(sid, 12, MUTED)]], space_after=2)
    x = Emu(int(x) + int(col_w) + Inches(0.05))

add_text(s, Inches(1.0), Inches(5.35), Inches(11), Inches(0.7),
         [[R("Supervised by", 11, MUTED, bold=True)],
          [R("Jaminur Rahman", 16, INK, bold=True),
           R("   ·   Lecturer, Department of CSE", 13, BODY)]],
         space_after=3)

add_text(s, Inches(1.0), Inches(6.72), Inches(11.5), Inches(0.6),
         [[R("Team Sherlocked", 13, WHITE, bold=True),
           R("    ·    Course CSE-3240 (Project I)    ·    "
             "Department of Computer Science & Engineering    ·    "
             "Leading University, Sylhet, Bangladesh", 12, RGBColor(0xCF, 0xD3, 0xDA))]],
         anchor=MSO_ANCHOR.MIDDLE)


# ============================================================================
# SLIDE 2 — PRESENTATION OUTLINE
# ============================================================================
s = new_slide()
header(s, "Presentation Outline")
footer(s, 2)

outline = [
    ("01", "Project Overview & Goal", "Idea, problem statement, and objectives"),
    ("02", "Functional Requirements", "What the system does + use-case model"),
    ("03", "Non-Functional Requirements", "Environment, frameworks, and quality goals"),
    ("04", "Technical Diagrams", "ER, DFD (Level 0 & 1), and activity diagrams"),
    ("05", "Project Demonstration", "Key features and the system in action"),
    ("06", "Discussion & Conclusion", "Outcomes, limitations, and future scope"),
    ("07", "Live Demo", "Functional walkthrough on a real device"),
    ("08", "Q & A", "Questions from the panel"),
]
top = 1.85
colx = [0.85, 6.95]
for i, (num, title, desc) in enumerate(outline):
    cx = colx[i // 4]
    ty = top + (i % 4) * 1.18
    add_text(s, Inches(cx), Inches(ty), Inches(0.85), Inches(0.8),
             [[R(num, 30, ORANGE, bold=True, font=FONT_LIGHT)]])
    add_text(s, Inches(cx + 0.95), Inches(ty + 0.02), Inches(4.9), Inches(1.0),
             [[R(title, 16, INK, bold=True)],
              [R(desc, 12, MUTED)]], space_after=2, line_spacing=1.05)


# ============================================================================
# SLIDE 3 — TEAM OVERVIEW
# ============================================================================
s = new_slide()
header(s, "Team Overview", kicker="Team Sherlocked")
footer(s, 3)

members = [
    ("Fahmid Alam", "0182320012101309",
     ["System architecture & shared infrastructure",
      "Authentication (Google OAuth + email)",
      "Admin module & timetable engine (OR-Tools CP-SAT)",
      "Database design, build & deployment"]),
    ("Swadheen Islam Robi", "0182320012101278",
     ["Student & teacher dashboards (frontend)",
      "Routine and resources hub",
      "Teacher Manage Classes module",
      "Find Teacher & Room Availability (real-time)"]),
    ("Shahriar Rashid Ratul", "0182320012101276",
     ["Notifications & push delivery (FCM)",
      "User profile module",
      "Quality assurance & data seeding",
      "Documentation & demo material"]),
]
cw = 3.95
gap = 0.18
x = 0.7
for name, sid, tasks in members:
    card = add_rect(s, Inches(x), Inches(1.75), Inches(cw), Inches(4.7),
                    WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
    card.line.color.rgb = LINE
    card.line.width = Pt(1)
    add_rect(s, Inches(x), Inches(1.75), Inches(cw), Inches(0.12), ORANGE,
             MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(x + 0.3), Inches(2.1), Inches(cw - 0.6), Inches(0.85),
             [[R(name, 16, INK, bold=True)],
              [R(sid, 11.5, MUTED)]], space_after=2)
    bullets(s, Inches(x + 0.3), Inches(3.0), Inches(cw - 0.55), Inches(3.3),
            [(t, 0) for t in tasks], size=12.5, gap=9)
    x += cw + gap


# ============================================================================
# SLIDE 4 — PROJECT OVERVIEW & GOAL
# ============================================================================
s = new_slide()
header(s, "Project Overview & Goal", kicker="Introduction")
footer(s, 4)

add_text(s, Inches(0.7), Inches(1.72), Inches(7.2), Inches(0.32),
         [[R("Description of the Project", 14, INK, bold=True)]])
add_text(s, Inches(0.7), Inches(2.12), Inches(7.2), Inches(0.9),
         [[R("UniVerse is an Android campus-companion app for the CSE department of "
             "Leading University. It automatically generates a clash-free department "
             "timetable and delivers real-time campus services to students, teachers, "
             "and admins from a single app.", 12.5, BODY)]], line_spacing=1.1)

add_text(s, Inches(0.7), Inches(3.35), Inches(7.2), Inches(0.32),
         [[R("Problem Statement & Motivation", 14, INK, bold=True)]])
bullets(s, Inches(0.7), Inches(3.74), Inches(7.2), Inches(1.4), [
    ("Manual timetable creation is slow, error-prone, and must satisfy many "
     "constraints (rooms, labs, teacher day-offs, Friday rule).", 0),
    ("Students and teachers lack one live source for today's classes, "
     "cancellations, free rooms, and teacher locations.", 0),
], size=12.5, gap=8)

add_text(s, Inches(0.7), Inches(5.25), Inches(7.2), Inches(0.32),
         [[R("Main Objectives & Goals", 14, INK, bold=True)]])
bullets(s, Inches(0.7), Inches(5.64), Inches(7.2), Inches(1.2), [
    ("Auto-generate a conflict-free timetable via constraint optimization.", 0),
    ("One role-aware app with real-time routine, rooms, and notifications.", 0),
], size=12.5, gap=8)

placeholder(s, Inches(8.2), Inches(1.85), Inches(4.5), Inches(4.85),
            "[ App Hero / Screenshots ]",
            "Add app mock-ups or a feature collage here")


# ============================================================================
# SLIDE 5 — FUNCTIONAL REQUIREMENTS
# ============================================================================
s = new_slide()
header(s, "Functional Requirements", kicker="Requirements")
footer(s, 5)

add_text(s, Inches(0.7), Inches(1.62), Inches(11.9), Inches(0.3),
         [[R("What the system must do, grouped by actor:", 12.5, BODY)]])

left_groups = [
    ("Authentication & Access (all users)", [
        "Sign up / sign in with Google or email",
        "Role-based access: student, teacher, admin",
        "Admin access gated by a whitelist"]),
    ("Student", [
        "View personal routine & dashboard (live / next class)",
        "Browse resources by semester folder",
        "Find a teacher and check room availability",
        "Receive in-app and push notifications"]),
]
right_groups = [
    ("Teacher", [
        "View personal teaching schedule",
        "Cancel a class occurrence (auto-notify students)",
        "Post a notice or room-change update"]),
    ("Admin", [
        "Auto-generate the department timetable (OR-Tools)",
        "Manage rooms, faculty, and settings; publish routine",
        "Broadcast notifications; manage resources & users"]),
]
grouped(s, Inches(0.7), Inches(2.1), Inches(5.9), Inches(4.7), left_groups)
grouped(s, Inches(6.9), Inches(2.1), Inches(5.7), Inches(4.7), right_groups)


# ============================================================================
# SLIDE 6 — USE CASE DIAGRAM
# ============================================================================
s = new_slide()
header(s, "Use Case Diagram", kicker="Functional Requirements")
footer(s, 6)
caption(s, "Actors (Student · Teacher · Admin) and their interactions with the system.")
placeholder(s, Inches(0.7), Inches(2.05), Inches(11.93), Inches(4.65),
            "[ Use Case Diagram ]",
            "Insert use-case-diagram.png here")


# ============================================================================
# SLIDE 7 — NON-FUNCTIONAL REQUIREMENTS
# ============================================================================
s = new_slide()
header(s, "Non-Functional Requirements", kicker="Requirements")
footer(s, 7)

grouped(s, Inches(0.7), Inches(1.85), Inches(5.9), Inches(4.9), [
    ("Operating System Environment", [
        "Client: Android (Flutter) phones",
        "Backend: Supabase (Postgres, Auth, Storage, Realtime)",
        "Engine: FastAPI + OR-Tools on Render (cloud)",
        "Push: Firebase Cloud Messaging"]),
    ("Frameworks & Tools", [
        "Flutter / Dart · GoRouter · Hive",
        "Supabase Flutter SDK · Firebase Messaging",
        "Python · FastAPI · Google OR-Tools (CP-SAT)",
        "Phosphor icons · Google Fonts (Inter)"]),
])
grouped(s, Inches(6.9), Inches(1.85), Inches(5.7), Inches(4.9), [
    ("Quality Attributes", [
        "Performance: clash-free schedule within a time limit",
        "Security: row-level security, OAuth PKCE, key isolation",
        "Reliability: real-time sync + push delivery",
        "Usability: dark theme, role-aware navigation",
        "Maintainability: layered Screen-Controller-Service",
        "Portability: cloud-hosted, stateless engine"]),
])


# ============================================================================
# SLIDES 8-11 — TECHNICAL DIAGRAMS
# ============================================================================
tech = [
    (8,  "Entity-Relationship Diagram",
     "Core tables: profiles · routines · cancellations · notifications · "
     "resources · timetable_* config · device_tokens.",
     "[ ER Diagram ]", "Insert er-diagram.png here"),
    (9,  "Data Flow Diagram - Level 0",
     "Context view: users interact with UniVerse, which talks to Supabase, "
     "the timetable engine, and FCM.",
     "[ DFD Level 0 (Context) ]", "Insert dfd-level-0.png here"),
    (10, "Data Flow Diagram - Level 1",
     "Decomposed processes: auth, routine, resources, notifications, and "
     "timetable generation.",
     "[ DFD Level 1 ]", "Insert dfd-level-1.png here"),
    (11, "Activity Diagram",
     "Flow of the automatic timetable generation: upload -> solve -> "
     "validate -> publish.",
     "[ Activity Diagram ]", "Insert activity-diagram.png here"),
]
for n, title, cap, ph_label, ph_sub in tech:
    s = new_slide()
    header(s, title, kicker="Technical Diagrams")
    footer(s, n)
    caption(s, cap)
    placeholder(s, Inches(0.7), Inches(2.05), Inches(11.93), Inches(4.65),
                ph_label, ph_sub)


# ============================================================================
# SLIDE 12 — PROJECT DEMONSTRATION
# ============================================================================
s = new_slide()
header(s, "Project Demonstration", kicker="Demonstration")
footer(s, 12)
caption(s, "A concise (~1 minute) overview highlighting the key features in action.")
placeholder(s, Inches(0.7), Inches(2.05), Inches(11.93), Inches(4.65),
            "[ Demo Video / Screenshot Slideshow ]",
            "Embed the 1-minute walkthrough or a screenshot strip here")


# ============================================================================
# SLIDE 13 — DISCUSSION / CONCLUSION
# ============================================================================
s = new_slide()
header(s, "Discussion & Conclusion", kicker="Conclusion")
footer(s, 13)

add_text(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(0.5),
         [[R("Outcome:  ", 13, ORANGE, bold=True),
           R("a working app that auto-generates a clash-free department timetable "
             "and serves it live to students, teachers, and admins.", 13, BODY)]])

cols = [
    ("Limitations", [
        "Free cloud tier sleeps - first request has a cold-start delay.",
        "Scoped to one department's distribution format.",
        "AI assistant (RAG) intentionally descoped for now.",
    ]),
    ("Restrictions", [
        "Android only - no iOS build yet.",
        "Release APK uses debug signing (sideload for demo).",
        "Depends on the standard course-distribution Excel file.",
    ]),
    ("Future Scope", [
        "AI campus assistant (RAG over documents).",
        "Multi-department & cross-faculty scheduling.",
        "iOS build and student-side cancellation grid.",
        "Solver tuning with student preferences.",
    ]),
]
x = 0.7
cw = 3.95
for title, items in cols:
    add_rect(s, Inches(x), Inches(2.35), Inches(cw), Pt(3.2), ORANGE)
    add_text(s, Inches(x), Inches(2.52), Inches(cw), Inches(0.4),
             [[R(title, 15, INK, bold=True)]])
    bullets(s, Inches(x), Inches(3.05), Inches(cw - 0.1), Inches(3.6),
            [(t, 0) for t in items], size=12.5, gap=10)
    x += cw + 0.18


# ============================================================================
# SLIDE 14 — DEMO
# ============================================================================
s = new_slide()
header(s, "Demo")
footer(s, 14)
placeholder(s, Inches(0.7), Inches(1.85), Inches(11.93), Inches(4.85),
            "[ Live Functional Demo ]",
            "Run the app live on a device / emulator here")


# ----------------------------------------------------------------------------
out = "docs/UniVerse_Project_Defense.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
