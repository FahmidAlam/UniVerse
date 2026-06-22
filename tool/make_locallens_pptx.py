"""
Generate a third-year PROJECT defense slide TEMPLATE (.pptx) for "LocalLens"
(a local-restaurant discovery, rating & review app).

Follows the department "Third Year Defense - Project Slide Show Template".
This is a FILL-IN TEMPLATE: branding + structure are done; bracketed [ ... ]
text and dashed boxes are placeholders for the team to complete.

Run:  python tool/make_locallens_pptx.py
Out:  docs/LocalLens_Project_Defense_Template.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ----------------------------------------------------------------------------
# Palette (LocalLens brand: gold on near-black)
# ----------------------------------------------------------------------------
GOLD        = RGBColor(0xC8, 0x91, 0x1B)   # rich amber accent on white
GOLD_BRIGHT = RGBColor(0xE6, 0xB0, 0x3A)   # brighter gold on dark band
INK         = RGBColor(0x1A, 0x1A, 0x1C)   # near-black headings
SLATE       = RGBColor(0x15, 0x15, 0x18)   # near-black band (app vibe)
BODY        = RGBColor(0x33, 0x37, 0x3D)   # body text
MUTED       = RGBColor(0x7A, 0x7E, 0x85)   # captions / placeholders
LINE        = RGBColor(0xDD, 0xDF, 0xE3)   # hairlines
PH_BG       = RGBColor(0xF5, 0xF4, 0xEE)   # placeholder fill (warm)
PH_BORDER   = RGBColor(0xC9, 0xC2, 0xAE)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

FONT       = "Calibri"
FONT_LIGHT = "Calibri Light"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------------
def add_rect(slide, l, t, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def add_text(slide, l, t, w, h, runs, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=0, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for (text, size, color, bold, font, italic) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = font
    return tb


def R(text, size, color=BODY, bold=False, font=FONT, italic=False):
    return (text, size, color, bold, font, italic)


def bg(slide, color=WHITE):
    add_rect(slide, 0, 0, EMU_W, EMU_H, color)


def footer(slide, n):
    add_rect(slide, Inches(0.6), Inches(7.02), Inches(12.13), Pt(0.75), LINE)
    add_text(slide, Inches(0.6), Inches(7.08), Inches(10), Inches(0.3),
             [[R("LocalLens  ·  [ Team Name ], [ Department ], [ University ]",
                 9, MUTED)]], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(11.9), Inches(7.08), Inches(0.83), Inches(0.3),
             [[R(str(n), 9, MUTED, bold=True)]],
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def header(slide, title, kicker=None):
    title_top = Inches(0.88) if kicker else Inches(0.55)
    rule_top = Inches(1.5) if kicker else Inches(1.18)
    if kicker:
        add_text(slide, Inches(0.7), Inches(0.55), Inches(11), Inches(0.3),
                 [[R(kicker.upper(), 11, GOLD, bold=True)]])
    add_text(slide, Inches(0.7), title_top, Inches(12), Inches(0.7),
             [[R(title, 28, INK, bold=True, font=FONT_LIGHT)]])
    add_rect(slide, Inches(0.72), rule_top, Inches(0.62), Pt(3.2), GOLD)


def bullets(slide, l, t, w, h, items, *, size=15, gap=10, sub_size=12.5,
            color=BODY, italic=False, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    for text, lvl in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = 1.06
        p.space_after = Pt(gap if lvl == 0 else gap * 0.45)
        p.space_before = Pt(0)
        if lvl == 0:
            dot = p.add_run()
            dot.text = "●  "
            dot.font.size = Pt(9)
            dot.font.color.rgb = GOLD
            dot.font.name = FONT
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.italic = italic
            r.font.name = FONT
        else:
            p.level = 1
            dash = p.add_run()
            dash.text = "–  "
            dash.font.size = Pt(sub_size)
            dash.font.color.rgb = MUTED
            dash.font.name = FONT
            r = p.add_run()
            r.text = text
            r.font.size = Pt(sub_size)
            r.font.color.rgb = MUTED
            r.font.italic = italic
            r.font.name = FONT
    return tb


def grouped(slide, l, t, w, h, groups, *, hsize=13.5, bsize=12, gap=5,
            group_gap=11):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    for head, items in groups:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.space_before = Pt(0 if first else group_gap)
        p.space_after = Pt(4)
        p.line_spacing = 1.0
        first = False
        r = p.add_run()
        r.text = head
        r.font.size = Pt(hsize)
        r.font.bold = True
        r.font.color.rgb = INK
        r.font.name = FONT
        for it, is_ph in items:
            pp = tf.add_paragraph()
            pp.space_before = Pt(0)
            pp.space_after = Pt(gap)
            pp.line_spacing = 1.04
            dot = pp.add_run()
            dot.text = "●  "
            dot.font.size = Pt(8)
            dot.font.color.rgb = GOLD
            dot.font.name = FONT
            rr = pp.add_run()
            rr.text = it
            rr.font.size = Pt(bsize)
            rr.font.color.rgb = MUTED if is_ph else BODY
            rr.font.italic = is_ph
            rr.font.name = FONT
    return tb


def placeholder(slide, l, t, w, h, label, sub=""):
    box = add_rect(slide, l, t, w, h, PH_BG, MSO_SHAPE.ROUNDED_RECTANGLE)
    box.line.color.rgb = PH_BORDER
    box.line.width = Pt(1.25)
    box.line.dash_style = 2
    runs = [[R(label, 14, MUTED, bold=True)]]
    if sub:
        runs.append([R(sub, 11, MUTED, italic=True)])
    add_text(slide, l, t, w, h, runs,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=3)


def caption(slide, text):
    add_text(slide, Inches(0.7), Inches(1.62), Inches(11.9), Inches(0.3),
             [[R(text, 12.5, BODY)]])


def new_slide():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    return s


# items helpers: (text, is_placeholder)
def F(t):   # real / inferred content
    return (t, False)


def P(t):   # placeholder content (muted italic)
    return (t, True)


# ============================================================================
# SLIDE 1 — PROJECT TITLE
# ============================================================================
s = new_slide()
add_rect(s, 0, 0, Inches(0.28), EMU_H, GOLD)
add_rect(s, 0, Inches(6.55), EMU_W, Inches(0.95), SLATE)

add_text(s, Inches(1.0), Inches(0.7), Inches(11), Inches(0.35),
         [[R("THIRD YEAR PROJECT  ·  DEFENSE PRESENTATION", 12, GOLD, bold=True)]])

add_text(s, Inches(0.95), Inches(1.55), Inches(11.4), Inches(1.4),
         [[R("LocalLens", 46, INK, bold=True, font=FONT_LIGHT)]], line_spacing=1.0)

add_text(s, Inches(1.0), Inches(2.75), Inches(11.2), Inches(0.9),
         [[R("Discover, Rate & Review Local Restaurants — a gamified food-discovery app",
             18, BODY)]], line_spacing=1.1)

add_rect(s, Inches(1.02), Inches(3.7), Inches(0.85), Pt(3.5), GOLD)

# team placeholders
col_w = Inches(3.75)
x = Inches(1.0)
for i in range(3):
    add_text(s, x, Inches(4.15), col_w, Inches(0.6),
             [[R(f"[ Team Member {i+1} ]", 15, INK, bold=True)],
              [R("[ Student ID ]", 12, MUTED, italic=True)]], space_after=2)
    x = Emu(int(x) + int(col_w) + Inches(0.05))

add_text(s, Inches(1.0), Inches(5.25), Inches(11), Inches(0.7),
         [[R("Supervised by", 11, MUTED, bold=True)],
          [R("[ Supervisor Name ]", 16, INK, bold=True),
           R("   ·   [ Designation, Department ]", 13, MUTED, italic=True)]],
         space_after=3)

add_text(s, Inches(1.0), Inches(6.72), Inches(11.5), Inches(0.6),
         [[R("[ Team Name ]", 13, WHITE, bold=True),
           R("    ·    Course [ Code ]    ·    "
             "[ Department ]    ·    [ University ]", 12, RGBColor(0xCF, 0xCB, 0xBE))]],
         anchor=MSO_ANCHOR.MIDDLE)


# ============================================================================
# SLIDE 2 — PRESENTATION OUTLINE
# ============================================================================
s = new_slide()
header(s, "Presentation Outline")
footer(s, 2)
outline = [
    ("01", "Project Overview & Goal", "Idea, problem statement, and objectives"),
    ("02", "Functional Requirements", "What the system does + use-case model"),
    ("03", "Non-Functional Requirements", "Environment, frameworks, and quality goals"),
    ("04", "Technical Diagrams", "ER, DFD (Level 0 & 1), and class diagram"),
    ("05", "Project Demonstration", "Key features and the system in action"),
    ("06", "Discussion & Conclusion", "Outcomes, limitations, and future scope"),
    ("07", "Live Demo", "Functional walkthrough on a real device"),
    ("08", "Q & A", "Questions from the panel"),
]
top = 1.85
colx = [0.85, 6.95]
for i, (num, title, desc) in enumerate(outline):
    cx = colx[i // 4]
    ty = top + (i % 4) * 1.18
    add_text(s, Inches(cx), Inches(ty), Inches(0.85), Inches(0.8),
             [[R(num, 30, GOLD, bold=True, font=FONT_LIGHT)]])
    add_text(s, Inches(cx + 0.95), Inches(ty + 0.02), Inches(4.9), Inches(1.0),
             [[R(title, 16, INK, bold=True)],
              [R(desc, 12, MUTED)]], space_after=2, line_spacing=1.05)


# ============================================================================
# SLIDE 3 — TEAM OVERVIEW
# ============================================================================
s = new_slide()
header(s, "Team Overview", kicker="The Team")
footer(s, 3)
cw = 3.95
gap = 0.18
x = 0.7
for i in range(3):
    card = add_rect(s, Inches(x), Inches(1.75), Inches(cw), Inches(4.7),
                    WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
    card.line.color.rgb = LINE
    card.line.width = Pt(1)
    add_rect(s, Inches(x), Inches(1.75), Inches(cw), Inches(0.12), GOLD,
             MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(x + 0.3), Inches(2.1), Inches(cw - 0.6), Inches(0.85),
             [[R(f"[ Team Member {i+1} ]", 16, INK, bold=True)],
              [R("[ Student ID ]", 11.5, MUTED, italic=True)]], space_after=2)
    bullets(s, Inches(x + 0.3), Inches(3.0), Inches(cw - 0.55), Inches(3.3),
            [("[ Contribution detail ]", 0)] * 4,
            size=12.5, gap=9, color=MUTED, italic=True)
    x += cw + gap


# ============================================================================
# SLIDE 4 — PROJECT OVERVIEW & GOAL
# ============================================================================
s = new_slide()
header(s, "Project Overview & Goal", kicker="Introduction")
footer(s, 4)

add_text(s, Inches(0.7), Inches(1.72), Inches(7.2), Inches(0.32),
         [[R("Description of the Project", 14, INK, bold=True)]])
add_text(s, Inches(0.7), Inches(2.12), Inches(7.2), Inches(0.95),
         [[R("LocalLens is a mobile app for discovering, rating, and reviewing "
             "local restaurants. It ranks the best dining spots and the most "
             "active reviewers through leaderboards, helping users find great "
             "places near them.  ", 12.5, BODY),
           R("[ edit to match your scope ]", 11, MUTED, italic=True)]],
         line_spacing=1.1)

add_text(s, Inches(0.7), Inches(3.45), Inches(7.2), Inches(0.32),
         [[R("Problem Statement & Motivation", 14, INK, bold=True)]])
bullets(s, Inches(0.7), Inches(3.84), Inches(7.2), Inches(1.4), [
    ("[ What problem does LocalLens solve? Why does it matter? ]", 0),
    ("[ Add background / motivation here ]", 0),
], size=12.5, gap=8, color=MUTED, italic=True)

add_text(s, Inches(0.7), Inches(5.3), Inches(7.2), Inches(0.32),
         [[R("Main Objectives & Goals", 14, INK, bold=True)]])
bullets(s, Inches(0.7), Inches(5.69), Inches(7.2), Inches(1.2), [
    ("[ Objective 1 ]", 0),
    ("[ Objective 2 ]", 0),
], size=12.5, gap=8, color=MUTED, italic=True)

placeholder(s, Inches(8.2), Inches(1.85), Inches(4.5), Inches(4.85),
            "[ App Hero / Screenshots ]",
            "Add LocalLens screenshots or a feature collage here")


# ============================================================================
# SLIDE 5 — FUNCTIONAL REQUIREMENTS  (inferred from screenshots; edit freely)
# ============================================================================
s = new_slide()
header(s, "Functional Requirements", kicker="Requirements")
footer(s, 5)
add_text(s, Inches(0.7), Inches(1.62), Inches(11.9), Inches(0.3),
         [[R("What the system must do (starter list from the app — edit as needed):",
             12.5, BODY)]])

left_groups = [
    ("Account & Profile", [
        F("Sign up / sign in"),
        F("Edit profile: name, bio, avatar, dark theme"),
        F("Account verification"),
        F("Earn reputation tiers (Explorer → Expert → Diamond → Platinum)")]),
    ("Discovery", [
        F("Search for dishes or places"),
        F("Browse the Top 10 restaurant leaderboard by location"),
        F("Find restaurants nearby (rating, cuisine, distance, status)"),
        F("View restaurants on a map")]),
]
right_groups = [
    ("Reviews & Gamification", [
        F("View restaurant details and ratings"),
        F("Write reviews and rate restaurants"),
        F("Earn points and receive “helpful” votes"),
        F("Appear on the Top 10 Critics leaderboard")]),
    ("Administration  (optional)", [
        P("[ Manage restaurant listings ]"),
        P("[ Verify users / moderate reviews ]")]),
]
grouped(s, Inches(0.7), Inches(2.1), Inches(5.9), Inches(4.7), left_groups)
grouped(s, Inches(6.9), Inches(2.1), Inches(5.7), Inches(4.7), right_groups)


# ============================================================================
# SLIDE 6 — USE CASE DIAGRAM
# ============================================================================
s = new_slide()
header(s, "Use Case Diagram", kicker="Functional Requirements")
footer(s, 6)
caption(s, "Actors (e.g., User / Reviewer · Admin) and their interactions with the system.")
placeholder(s, Inches(0.7), Inches(2.05), Inches(11.93), Inches(4.65),
            "[ Use Case Diagram ]", "Insert your use-case diagram here")


# ============================================================================
# SLIDE 7 — NON-FUNCTIONAL REQUIREMENTS
# ============================================================================
s = new_slide()
header(s, "Non-Functional Requirements", kicker="Requirements")
footer(s, 7)
grouped(s, Inches(0.7), Inches(1.85), Inches(5.9), Inches(4.9), [
    ("Operating System Environment", [
        P("[ Client platform: e.g., Android / iOS ]"),
        P("[ Backend: e.g., Firebase / Supabase / Node.js ]"),
        P("[ Location & maps service ]"),
        P("[ Image / cloud storage ]")]),
    ("Frameworks & Tools", [
        P("[ Framework: e.g., Flutter / React Native ]"),
        P("[ Database & SDKs ]"),
        P("[ Other libraries and tools ]")]),
])
grouped(s, Inches(6.9), Inches(1.85), Inches(5.7), Inches(4.9), [
    ("Quality Attributes", [
        F("Performance: fast search and nearby results"),
        F("Accuracy: reliable location and distance"),
        F("Security: secure authentication and user data"),
        F("Usability: clean dark UI, simple navigation"),
        F("Scalability: grows with more places and users"),
        F("Reliability: consistent ratings and sync")]),
])


# ============================================================================
# SLIDES 8-11 — TECHNICAL DIAGRAMS
# ============================================================================
tech = [
    (8,  "Entity-Relationship Diagram",
     "Main entities, e.g., users, restaurants, reviews, ratings, points/badges.",
     "[ ER Diagram ]", "Insert your ER diagram here"),
    (9,  "Data Flow Diagram - Level 0",
     "Context view: the user, the LocalLens system, and external services.",
     "[ DFD Level 0 (Context) ]", "Insert your DFD Level 0 here"),
    (10, "Data Flow Diagram - Level 1",
     "Decomposed processes: auth, discovery, reviews, leaderboards.",
     "[ DFD Level 1 ]", "Insert your DFD Level 1 here"),
    (11, "Class Diagram",
     "Core classes and their relationships.",
     "[ Class Diagram ]", "Insert your class diagram here"),
]
for n, title, cap, ph_label, ph_sub in tech:
    s = new_slide()
    header(s, title, kicker="Technical Diagrams")
    footer(s, n)
    caption(s, cap)
    placeholder(s, Inches(0.7), Inches(2.05), Inches(11.93), Inches(4.65),
                ph_label, ph_sub)


# ============================================================================
# SLIDE 12 — PROJECT DEMONSTRATION
# ============================================================================
s = new_slide()
header(s, "Project Demonstration", kicker="Demonstration")
footer(s, 12)
caption(s, "A concise (~1 minute) overview highlighting the key features in action.")
placeholder(s, Inches(0.7), Inches(2.05), Inches(11.93), Inches(4.65),
            "[ Demo Video / Screenshot Slideshow ]",
            "Embed the 1-minute walkthrough or a screenshot strip here")


# ============================================================================
# SLIDE 13 — DISCUSSION / CONCLUSION
# ============================================================================
s = new_slide()
header(s, "Discussion & Conclusion", kicker="Conclusion")
footer(s, 13)
add_text(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(0.5),
         [[R("Outcome:  ", 13, GOLD, bold=True),
           R("[ summarize what the project achieved in one line ]", 13, MUTED,
             italic=True)]])
cols = [
    ("Limitations", ["[ Limitation 1 ]", "[ Limitation 2 ]", "[ Limitation 3 ]"]),
    ("Restrictions", ["[ Restriction 1 ]", "[ Restriction 2 ]", "[ Restriction 3 ]"]),
    ("Future Scope", ["[ Future improvement 1 ]", "[ Future improvement 2 ]",
                      "[ Future improvement 3 ]"]),
]
x = 0.7
cw = 3.95
for title, items in cols:
    add_rect(s, Inches(x), Inches(2.35), Inches(cw), Pt(3.2), GOLD)
    add_text(s, Inches(x), Inches(2.52), Inches(cw), Inches(0.4),
             [[R(title, 15, INK, bold=True)]])
    bullets(s, Inches(x), Inches(3.05), Inches(cw - 0.1), Inches(3.6),
            [(t, 0) for t in items], size=12.5, gap=10, color=MUTED, italic=True)
    x += cw + 0.18


# ============================================================================
# SLIDE 14 — DEMO
# ============================================================================
s = new_slide()
header(s, "Demo")
footer(s, 14)
placeholder(s, Inches(0.7), Inches(1.85), Inches(11.93), Inches(4.85),
            "[ Live Functional Demo ]",
            "Run the app live on a device / emulator here")


out = "docs/LocalLens_Project_Defense_Template.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
